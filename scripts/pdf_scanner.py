#!/usr/bin/env python3
"""
Axiom — Universal PII scanner subprocess (Gemma-3-4B + Regex edition)
=======================================================================
Supports: PDF, DOCX, XLSX, PPTX
"""

import sys, json, re, logging, os, time

logging.basicConfig(stream=sys.stderr, level=logging.INFO,
                    format="[PDFScanner] %(levelname)s %(message)s")
log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

GEMMA_MODEL_ID = os.environ.get("AXIOM_GEMMA_MODEL", "google/gemma-3-4b-it")
CHUNK_SIZE     = int(os.environ.get("AXIOM_CHUNK_SIZE",    "3000"))   
CHUNK_OVERLAP  = int(os.environ.get("AXIOM_CHUNK_OVERLAP", "200"))
MAX_NEW_TOKENS = 1024

LABEL_SEVERITY = {
    "SSN":             "CRITICAL",
    "CREDIT_CARD":     "CRITICAL",
    "CC_METADATA":     "CRITICAL",
    "AWS_KEY":         "CRITICAL",
    "AWS_SECRET":      "CRITICAL",
    "JWT":             "CRITICAL",
    "PRIVATE_KEY":     "CRITICAL",
    "PASSWORD":        "CRITICAL",
    "API_KEY":         "CRITICAL",
    "PASSPORT":        "CRITICAL",
    "BANK_ACCOUNT":    "CRITICAL",
    "DRIVERS_LICENSE": "HIGH",
    "DOB":             "MEDIUM",
    "EMAIL":           "MEDIUM",
    "PHONE":           "MEDIUM",
    "NAME":            "LOW",
    "ADDRESS":         "LOW",
    "IP_ADDRESS":      "LOW",
}
DEFAULT_SEVERITY = "MEDIUM"
VALID_SEVERITIES = {"CRITICAL", "HIGH", "MEDIUM", "LOW"}

# ---------------------------------------------------------------------------
# Model state
# ---------------------------------------------------------------------------

_model     = None
_tokenizer = None
_device    = None

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def redact(text: str) -> str:
    return text[:4] + "****" if len(text) > 4 else "*" * len(text)

def normalize_severity(s: str) -> str:
    s = s.upper().strip()
    return s if s in VALID_SEVERITIES else DEFAULT_SEVERITY

def normalize_string(s: str) -> str:
    """Removes spaces, punctuation, and casing for bulletproof deduplication."""
    return re.sub(r'\W+', '', s).lower()

def chunk_text(text: str) -> list:
    step = max(1, CHUNK_SIZE - CHUNK_OVERLAP)
    chunks = []
    for i in range(0, len(text), step):
        c = text[i : i + CHUNK_SIZE]
        if c.strip():
            chunks.append(c)
    return chunks or [text]

# ---------------------------------------------------------------------------
# Stage 1 — Regex
# ---------------------------------------------------------------------------

_DOB_CONTEXT_RE = re.compile(
    r"\b(birth|dob|date[\s_\-]of[\s_\-]birth|born|birthdate|birthday)\b",
    re.IGNORECASE
)

_REGEX_PATTERNS = [
    ("SSN",         "CRITICAL", re.compile(r"\b\d{3}-\d{2}-\d{4}\b")),
    ("EMAIL",       "MEDIUM",   re.compile(r"\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b")),
    ("IP_ADDRESS",  "LOW",      re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b")),
    ("PHONE",       "MEDIUM",   re.compile(r"\b(?:\+1[\s\-]?)?\(?\d{3}\)?[\s\-]?\d{3}[\s\-]?\d{4}\b")),
    ("CREDIT_CARD", "CRITICAL", re.compile(r"\b(?:4\d{12}(?:\d{3})?|5[1-5]\d{14}|3[47]\d{13}|6(?:011|5\d{2})\d{12})\b")),
    ("CC_METADATA", "CRITICAL", re.compile(r"(?i)\b(?:cvv|cvc|exp|expiration)\b\s*[:=]?\s*[\d/]+")),
    ("AWS_KEY",     "CRITICAL", re.compile(r"\bAKIA[0-9A-Z]{16}\b")),
    ("AWS_SECRET",  "CRITICAL", re.compile(r"(?i)(?:secret|aws_secret_access_key).{0,20}[:=]\s*([A-Za-z0-9/+= \n\r]{35,45})")),
    ("PASSWORD",    "CRITICAL", re.compile(r"(?i)(?:password|passwd|pwd).{0,20}[:=]\s*([^\s]{5,})")),
    ("JWT",         "CRITICAL", re.compile(r"\beyJ[A-Za-z0-9_\-]{10,}\.[A-Za-z0-9_\-]{10,}\.[A-Za-z0-9_\-]{10,}\b")),
    ("PRIVATE_KEY", "CRITICAL", re.compile(r"-----BEGIN (?:RSA |EC |OPENSSH )?PRIVATE KEY-----")),
]

_DATE_RE = re.compile(r"\b(?:\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4}|\d{4}[\/\-\.]\d{1,2}[\/\-\.]\d{1,2})\b")

def regex_scan(text: str, page_num: int) -> list:
    results = []
    seen = set()

    for label, severity, pat in _REGEX_PATTERNS:
        for m in pat.finditer(text):
            val = m.group(1) if m.lastindex else m.group(0)
            val = re.sub(r'\s+', ' ', val).strip()
            if not val:
                continue

            norm_val = normalize_string(val)
            key = (label, norm_val)
            if key in seen:
                continue
            seen.add(key)
            
            log.info(f"  p{page_num} [regex/high] {label}: {redact(val)!r}")
            results.append({
                "label":      label,
                "severity":   severity,
                "redacted":   redact(val),
                "value":      val,          
                "page":       page_num,
                "confidence": "high",
                "source":     "regex",
            })

    for m in _DATE_RE.finditer(text):
        val = m.group(0).strip()
        norm_val = normalize_string(val)
        key = ("DOB", norm_val)
        
        if key in seen:
            continue
            
        start = max(0, m.start() - 120)
        end   = min(len(text), m.end() + 120)
        ctx   = text[start:end]
        
        if _DOB_CONTEXT_RE.search(ctx):
            seen.add(key)
            log.info(f"  p{page_num} [regex/high] DOB (context): {redact(val)!r}")
            results.append({
                "label":      "DOB",
                "severity":   "MEDIUM",
                "redacted":   redact(val),
                "value":      val,
                "page":       page_num,
                "confidence": "high",
                "source":     "regex",
            })

    return results

# ---------------------------------------------------------------------------
# Stage 2 — Gemma 4B
# ---------------------------------------------------------------------------

def load_model() -> None:
    global _model, _tokenizer, _device
    import torch
    from transformers import AutoTokenizer, AutoModelForCausalLM

    _device = ("mps" if torch.backends.mps.is_available() else "cuda" if torch.cuda.is_available() else "cpu")
    log.info(f"Loading {GEMMA_MODEL_ID} on {_device}...")
    _tokenizer = AutoTokenizer.from_pretrained(GEMMA_MODEL_ID)
    _model = AutoModelForCausalLM.from_pretrained(GEMMA_MODEL_ID, dtype=torch.bfloat16).to(_device)
    _model.eval()
    log.info(f"Gemma ready (device={_device})")

GEMMA_SYSTEM = (
    "You are a data-loss-prevention engine. Extract ONLY personally identifiable information (PII) "
    "and sensitive credentials from the TEXT. Return a JSON array — nothing else.\n\n"
    "Rules:\n"
    "- A bare date like 01/01/1990 is only DOB if the text explicitly says 'date of birth', 'DOB', 'born', or 'birthday' nearby.\n"
    "- Do NOT invent values. Only extract text that literally appears in the TEXT.\n"
    "- Do NOT repeat the same value twice.\n"
    "- If nothing sensitive is present return [].\n\n"
    "Labels (use exactly): NAME, ADDRESS, DOB, PASSPORT, DRIVERS_LICENSE, BANK_ACCOUNT, EMAIL, PHONE, SSN, CREDIT_CARD, IP_ADDRESS, AWS_KEY, JWT, PRIVATE_KEY, PASSWORD, API_KEY\n\n"
    "Each object: "
    '{"label":"...","value":"exact text","severity":"CRITICAL|HIGH|MEDIUM|LOW","confidence":"high|medium|low"}'
)

def _parse_json_array(text: str) -> list:
    clean = re.sub(r"```(?:json)?", "", text).strip().rstrip("`").strip()
    try:
        r = json.loads(clean)
        if isinstance(r, list): return r
    except json.JSONDecodeError: pass
    m = re.search(r"\[[\s\S]*\]", clean)
    if m:
        try:
            r = json.loads(m.group(0))
            if isinstance(r, list): return r
        except json.JSONDecodeError: pass
    return []

def run_gemma(text_chunk: str) -> tuple:
    import torch
    messages = [{"role": "user", "content": f"{GEMMA_SYSTEM}\n\nTEXT:\n{text_chunk}"}]
    prompt = _tokenizer.apply_chat_template(messages, tokenize=False, add_generation_prompt=True)
    inputs = _tokenizer(prompt, return_tensors="pt", truncation=True, max_length=4096).to(_device)

    with torch.inference_mode():
        output_ids = _model.generate(**inputs, max_new_tokens=MAX_NEW_TOKENS, do_sample=False, pad_token_id=_tokenizer.eos_token_id)

    input_len = inputs["input_ids"].shape[1]
    raw = _tokenizer.decode(output_ids[0][input_len:], skip_special_tokens=True).strip()
    return _parse_json_array(raw), raw

def gemma_scan(text: str, page_num: int, regex_values_normalized: set) -> tuple:
    chunks       = chunk_text(text)
    seen         = set(regex_values_normalized)
    dets         = []
    raw_count    = 0
    parse_errors = 0
    t0           = time.monotonic()

    log.info(f"  p{page_num}: {len(chunks)} chunk(s) → Gemma")

    for idx, chunk in enumerate(chunks):
        try:
            findings, raw_out = run_gemma(chunk)
        except Exception as e:
            parse_errors += 1
            continue

        if not findings:
            clean_raw = re.sub(r"```(?:json)?|```", "", raw_out).strip()
            is_empty = clean_raw in ("[]", "", "[ ]")
            if not is_empty: parse_errors += 1
            continue

        raw_count += len(findings)

        for f in findings:
            label      = str(f.get("label", "UNKNOWN")).upper().strip()
            value      = str(f.get("value", "")).strip()
            severity   = normalize_severity(str(f.get("severity", LABEL_SEVERITY.get(label, DEFAULT_SEVERITY))))
            confidence = str(f.get("confidence", "medium")).lower().strip()

            norm_val = normalize_string(value)
            if not norm_val or norm_val in seen:
                continue

            if label == "DOB":
                start = max(0, text.find(value) - 120)
                end   = min(len(text), text.find(value) + len(value) + 120)
                if not _DOB_CONTEXT_RE.search(text[start:end]):
                    continue

            seen.add(norm_val)
            log.info(f"  p{page_num} [gemma/{confidence}] {label}: {redact(value)!r}")
            dets.append({
                "label":      label,
                "severity":   severity,
                "redacted":   redact(value),
                "value":      value,
                "page":       page_num,
                "confidence": confidence,
                "source":     "gemma",
            })

    duration_ms = round((time.monotonic() - t0) * 1000)
    status = "parse_error" if (parse_errors > 0 and not dets) else ("empty" if not dets else "ok")

    page_log = {
        "page": page_num, "chunks": len(chunks),
        "raw_findings": raw_count, "kept": len(dets),
        "duration_ms": duration_ms, "status": status,
    }
    return dets, page_log

# ---------------------------------------------------------------------------
# Text extractors
# ---------------------------------------------------------------------------

def extract_pdf(path: str) -> list:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer
    pages = []
    for page_num, layout in enumerate(extract_pages(path), start=1):
        text = "".join(el.get_text() for el in layout if isinstance(el, LTTextContainer))
        pages.append((page_num, text))
    return pages

def extract_docx(path: str) -> list:
    from docx import Document
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += "\n" + cell.text
    return [(1, text)]

def extract_xlsx(path: str) -> list:
    import openpyxl
    # Do NOT use read_only=True — it causes data_only=True to return None for
    # formula-result cells in many real-world xlsx files, starving Gemma of text.
    # data_only=True still gives us the cached formula results when available.
    wb = openpyxl.load_workbook(path, data_only=True)
    pages = []
    for sheet_num, ws in enumerate(wb.worksheets, start=1):
        rows = []
        # Include the sheet title as context so Gemma understands the domain
        rows.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            row_text = "\t".join(cells)
            if row_text.strip():
                rows.append(row_text)
        if len(rows) > 1:  # more than just the sheet title
            pages.append((sheet_num, "\n".join(rows)))
    wb.close()
    return pages

def extract_pptx(path: str) -> list:
    from pptx import Presentation
    prs = Presentation(path)
    pages = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.extend(p.text for p in shape.text_frame.paragraphs)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text_frame.text if cell.text_frame else "")
        pages.append((slide_num, "\n".join(parts)))
    return pages

def extract_document(path: str) -> list:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf": return extract_pdf(path)
    elif ext in (".docx", ".doc"): return extract_docx(path)
    elif ext in (".xlsx", ".xls", ".csv"): return extract_xlsx(path)
    elif ext in (".pptx", ".ppt"): return extract_pptx(path)
    else: raise ValueError(f"Unsupported file type: {ext}")

# ---------------------------------------------------------------------------
# Full scan pipeline
# ---------------------------------------------------------------------------

def scan_document(path: str) -> tuple:
    log.info(f"Scanning: {path}")
    pages = extract_document(path)
    
    all_dets = []
    all_logs = []

    for page_num, text in pages:
        if not text.strip():
            all_logs.append({
                "page": page_num, "chunks": 0, "raw_findings": 0,
                "kept": 0, "duration_ms": 0, "status": "blank",
            })
            continue

        regex_dets = regex_scan(text, page_num)
        regex_values_normalized = {normalize_string(d["value"]) for d in regex_dets}

        gemma_dets, page_log = gemma_scan(text, page_num, regex_values_normalized)

        page_dets = []
        for d in regex_dets + gemma_dets:
            d["raw_value"] = d.pop("value", "")
            d.pop("source", None)
            page_dets.append(d)

        all_dets.extend(page_dets)
        all_logs.append(page_log)

    return all_dets, all_logs

# ---------------------------------------------------------------------------
# Main loop
# ---------------------------------------------------------------------------

def check_deps(ext: str) -> str | None:
    try: import pdfminer  # noqa
    except ImportError:
        if ext == ".pdf": return "pdfminer.six not installed"
    try: import docx  # noqa
    except ImportError:
        if ext in (".docx", ".doc"): return "python-docx not installed"
    try: import openpyxl  # noqa
    except ImportError:
        if ext in (".xlsx", ".xls"): return "openpyxl not installed"
    try: import pptx  # noqa
    except ImportError:
        if ext in (".pptx", ".ppt"): return "python-pptx not installed"
    return None

def main():
    try: import pdfminer  # noqa
    except ImportError:
        sys.stdout.write(json.dumps({
            "id": "init", "detections": [], "ready": False,
            "error": "pdfminer.six not installed — run: pip install pdfminer.six",
        }) + "\n")
        sys.stdout.flush()
        sys.exit(1)

    try: load_model()
    except Exception as e:
        sys.stdout.write(json.dumps({
            "id": "init", "detections": [], "ready": False,
            "error": f"Gemma failed to load: {e}",
        }) + "\n")
        sys.stdout.flush()
        sys.exit(1)

    sys.stdout.write(json.dumps({
        "id": "init", "detections": [], "error": None, "ready": True,
        "model": GEMMA_MODEL_ID, "device": _device, "chunk_size": CHUNK_SIZE,
    }) + "\n")
    sys.stdout.flush()

    for line in sys.stdin:
        line = line.strip()
        if not line: continue
        try: req = json.loads(line)
        except json.JSONDecodeError as e:
            sys.stdout.write(json.dumps({"id": "?", "detections": [], "gemma_log": [], "error": f"JSON: {e}"}) + "\n")
            sys.stdout.flush()
            continue

        req_id   = req.get("id", "?")
        pdf_path = req.get("pdf_path", "")
        
        ext = os.path.splitext(pdf_path)[1].lower()
        dep_err = check_deps(ext)
        if dep_err:
            sys.stdout.write(json.dumps({"id": req_id, "detections": [], "gemma_log": [], "error": dep_err}) + "\n")
            sys.stdout.flush()
            continue

        try:
            dets, gemma_log = scan_document(pdf_path)
            resp = {"id": req_id, "detections": dets, "gemma_log": gemma_log, "error": None}
        except Exception as e:
            resp = {"id": req_id, "detections": [], "gemma_log": [], "error": str(e)}

        sys.stdout.write(json.dumps(resp) + "\n")
        sys.stdout.flush()

if __name__ == "__main__":
    main()